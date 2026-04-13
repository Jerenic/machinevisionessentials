"""
Erzeugt presentation/MachineVision_Gruppe8.pptx

Ausführen (im Projektroot):
  uv sync --extra slides
  uv run --extra slides python presentation/build_presentation.py
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt


def add_title_slide(prs: Presentation, title: str, subtitle: str) -> None:
    layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = title
    if slide.placeholders[1]:
        slide.placeholders[1].text = subtitle


def add_bullet_slide(prs: Presentation, title: str, bullets: list[str]) -> None:
    layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = title
    body = slide.placeholders[1]
    tf = body.text_frame
    tf.clear()
    for i, line in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.level = 0
        p.font.size = Pt(20)


def main() -> None:
    root = Path(__file__).resolve().parent.parent
    out = root / "presentation" / "MachineVision_Gruppe8.pptx"

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    add_title_slide(
        prs,
        "Machine Vision: Candy-Erkennung",
        "Förderband · YOLO11n · Label Studio\n\nGruppe 8\n"
        "Nicolas Haselbacher (PKZ 2410859016) · Jeremy Heißenberger (PKZ 2410859026)\n"
        "Hochschule Burgenland · MachineVisionEssentials",
    )

    add_bullet_slide(
        prs,
        "Aufgabe",
        [
            "Automatische Erkennung verschiedener Süßigkeitensorten auf einem Förderband",
            "Kamera von oben · Klassifikation + Zählung sichtbarer Stücke",
            "Umsetzung: Objekterkennung mit YOLO (Ultralytics), Training auf eigenem Datensatz",
        ],
    )

    add_bullet_slide(
        prs,
        "Pipeline (Überblick)",
        [
            "Fotos (Smartphone) → Konvertierung HEIC/JPEG",
            "Annotation in Label Studio (Bounding Boxes, 8 Klassen)",
            "Export YOLO · split_dataset.py → train / val / test + data.yaml",
            "Augmentation → ≥ 200 Trainingsbilder",
            "Training YOLO11n · Inferenz Test & Webcam",
        ],
    )

    add_bullet_slide(
        prs,
        "Klassen (inkl. Erweiterung)",
        [
            "Final trainierte Klassen: Bounty, Galaxy, GalaxyCaramel, Maltesers, Mars, MilkyWay, Snickers, Twix",
            "Erweiterung gegenüber einer generischen Candy-Erkennung: konkrete Sorten als eigene Klassen",
            "Wichtig für Vortrag: Welche Klassen haben wir bewusst getrennt? (z. B. Galaxy vs. GalaxyCaramel)",
        ],
    )

    add_bullet_slide(
        prs,
        "Training",
        [
            "Modell: YOLO11n (Transfer Learning)",
            "Mehrere Läufe: u. a. 50 Epochen vs. 100 Epochen",
            "Ergebnis: 100 Epochen → deutlich stabilere, bessere Erkennung",
            "Artefakte: runs/detect/ · z. B. candy_train2 (100 Epochen)",
        ],
    )

    add_bullet_slide(
        prs,
        "Ergebnisse",
        [
            "Gesamt: sehr gute Erkennung nach 100 Epochen (Val-Metriken, Tests)",
            "Schwachstellen: vor allem Mars und Milky Way",
            "Vermutung: zu wenige bzw. schwer unterscheidbare Trainingsbeispiele für diese Sorten",
            "Webcam: Nahaufnahmen oft korrekt; Hintergrund gelegentlich Fehlalarme (Domain Shift)",
        ],
    )

    add_bullet_slide(
        prs,
        "Herausforderungen",
        [
            "HEIC → JPG · Import in Label Studio",
            "Zeitaufwand: viele Bounding Boxes annotieren",
            "Windows: DataLoader workers=0 (Speicher / Auslagerungsdatei)",
            "Optional: mehr Fotos für Mars & Milky Way · höhere conf-Schwelle Webcam",
        ],
    )

    add_bullet_slide(
        prs,
        "Was bei uns nicht funktioniert hat",
        [
            "Mit 50 Epochen war die Erkennung deutlich schlechter als mit 100 Epochen",
            "Webcam: False Positives auf Nicht-Candy-Objekten (z. B. Gesicht/Hintergrund)",
            "Mars und MilkyWay bleiben schwächer (Verwechslungen/unsichere Scores)",
            "Training unter Windows anfangs instabil (WinError 1455, OOM in DataLoader-Workern)",
        ],
    )

    add_bullet_slide(
        prs,
        "Lösungen & Verständnisprobleme",
        [
            "Lösung: 100 Epochen + Augmentation + Workers=0 + Batch reduziert → stabiles Training",
            "Lösung: Klassen-/Datenfokus auf Problemklassen (Mars, MilkyWay) als nächster Schritt",
            "Verständnisproblem: README-Split-Angabe war missverständlich (80/10/10 richtig als train/val/test)",
            "Verständnisproblem: Detect vs. OBB/Keypoints sind unterschiedliche Tasks mit anderen Labels",
        ],
    )

    add_bullet_slide(
        prs,
        "Insights (Takeaways)",
        [
            "Mehr Epochen helfen nur mit passenden Daten – Datenqualität > reine Trainingsdauer",
            "Klassenbalance und harte Negativbeispiele reduzieren False Positives stark",
            "Webcam-Tests zeigen Domain-Shift früh und sind extrem wertvoll für Realitätscheck",
            "Für die Demo: Erfolge klar zeigen, Schwächen transparent benennen und nächste Schritte ableiten",
        ],
    )

    add_bullet_slide(
        prs,
        "Fazit",
        [
            "Pipeline von Rohfoto bis Live-Webcam funktioniert",
            "Mehr Epochen (100) lohnt sich klar gegenüber 50",
            "Klassenbalance und Domain (Webcam vs. Tisch) sind die nächsten Hebel",
            "Vielen Dank für die Aufmerksamkeit – Fragen?",
        ],
    )

    out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out))
    print(f"Gespeichert: {out}")


if __name__ == "__main__":
    main()
